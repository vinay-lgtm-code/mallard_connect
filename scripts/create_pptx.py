#!/usr/bin/env python3
"""Generate Mallard Connect UX Presentation (.pptx)"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand colors
TEAL_DARK = RGBColor(0x1A, 0x56, 0x53)
TEAL_MID = RGBColor(0x2D, 0x7D, 0x79)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xDC, 0x26, 0x26)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)
AMBER_LIGHT = RGBColor(0xFE, 0xF3, 0xC7)
GREEN_LIGHT = RGBColor(0xF0, 0xFD, 0xF4)
BLUE_LIGHT = RGBColor(0xEF, 0xF6, 0xFF)
PURPLE_LIGHT = RGBColor(0xFA, 0xF5, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, corner_radius=None):
    if corner_radius:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        # Rounded corners
        shape.adjustments[0] = corner_radius
    else:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=GRAY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
    return txBox


def add_screen_placeholder(slide, left, top, width, height, label, device="desktop"):
    """Add a mockup frame placeholder with label"""
    # Outer device frame
    if device == "mobile":
        shape = add_shape(slide, left, top, width, height, LIGHT_GRAY, 0.08)
        shape.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
        shape.line.width = Pt(2)
        # Status bar
        add_shape(slide, left + Inches(0.1), top + Inches(0.1), width - Inches(0.2), Inches(0.2), WHITE)
        add_text_box(slide, left + Inches(0.15), top + Inches(0.08), Inches(0.5), Inches(0.25),
                     "9:41", font_size=8, color=BLACK, bold=True)
    else:
        shape = add_shape(slide, left, top, width, height, LIGHT_GRAY, 0.02)
        shape.line.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
        shape.line.width = Pt(2)
        # Title bar
        add_shape(slide, left + Inches(0.05), top + Inches(0.05), width - Inches(0.1), Inches(0.25), WHITE)

    # Screen label
    add_text_box(slide, left, top + height + Inches(0.1), width, Inches(0.3),
                 label, font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, TEAL_DARK)

# Logo area
add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
             "MALLARD CONNECT", font_size=52, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(2.7), Inches(11), Inches(0.8),
             "Lead Nurturing & Follow-up System", font_size=28, color=AMBER)
add_text_box(slide, Inches(1), Inches(3.7), Inches(8), Inches(1),
             "A lightweight pre-CRM that ensures no prospect falls through the cracks.\nWorks alongside the Mortgage Advice Bureau Platform system.",
             font_size=16, color=WHITE)

# Accent bar
add_shape(slide, Inches(1), Inches(3.4), Inches(3), Inches(0.04), AMBER)

# Footer
add_text_box(slide, Inches(1), Inches(6.2), Inches(5), Inches(0.4),
             "Mallard Mortgages, Sheffield UK", font_size=14, color=RGBColor(0x7D, 0xB3, 0xB0))
add_text_box(slide, Inches(8), Inches(6.2), Inches(4.3), Inches(0.4),
             "Built by Storyboard Digital", font_size=14, color=RGBColor(0x7D, 0xB3, 0xB0),
             alignment=PP_ALIGN.RIGHT)


# ============================================================
# SLIDE 2: The Problem
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.7),
             "The Problem", font_size=36, color=TEAL_DARK, bold=True)
add_shape(slide, Inches(1), Inches(1.3), Inches(2), Inches(0.04), AMBER)

add_text_box(slide, Inches(1), Inches(1.8), Inches(5.5), Inches(1.5),
             "Prospects who enquire but aren't ready to proceed are falling through the cracks.",
             font_size=22, color=BLACK)

# Problem cards
problems = [
    ("No structured follow-up", "Advisers note \"contact in January\" but there's no system to make it happen.", RED),
    ("No pipeline visibility", "Della can't see at a glance where each prospect sits or what the next action should be.", AMBER),
    ("CRM is overkill (for now)", "A full CRM is too expensive and complex. They need something that works with what they have.", BLUE),
]
for i, (title, desc, color) in enumerate(problems):
    y = Inches(3.3) + Inches(i * 1.2)
    bar = add_shape(slide, Inches(1), y, Inches(0.06), Inches(0.8), color)
    add_text_box(slide, Inches(1.4), y, Inches(4.5), Inches(0.35), title, font_size=16, color=BLACK, bold=True)
    add_text_box(slide, Inches(1.4), y + Inches(0.35), Inches(4.5), Inches(0.5), desc, font_size=13, color=GRAY)

# Right side - solution box
sol_box = add_shape(slide, Inches(7), Inches(1.8), Inches(5.3), Inches(4.8), RGBColor(0xF0, 0xFD, 0xF4), 0.03)
sol_box.line.color.rgb = GREEN
sol_box.line.width = Pt(2)
add_text_box(slide, Inches(7.5), Inches(2.1), Inches(4.3), Inches(0.5),
             "The Solution: Mallard Connect", font_size=20, color=RGBColor(0x16, 0x65, 0x34), bold=True)
solutions = [
    "Structured follow-up reminders with email notifications to up to 3 recipients",
    "Visual pipeline board showing every prospect's status at a glance",
    "Daily 7am reminder emails so no lead is ever forgotten",
    "MAB Platform CSV import with automatic deduplication",
    "Mobile-optimized for salespeople in the field",
    "Real-time team dashboard for Della",
]
add_bullet_list(slide, Inches(7.5), Inches(2.7), Inches(4.3), Inches(3.5),
                [f"  {s}" for s in solutions], font_size=13, color=RGBColor(0x16, 0x65, 0x34))


# ============================================================
# SLIDE 3: Personas
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.7),
             "Three Personas, One System", font_size=36, color=TEAL_DARK, bold=True)
add_shape(slide, Inches(1), Inches(1.3), Inches(2), Inches(0.04), AMBER)

personas = [
    ("Della Mallard", "Owner / Manager", TEAL_DARK, [
        "Sees everything at a glance",
        "Pipeline health & team performance",
        "Assigns and rebalances leads",
        "Imports from MAB Platform",
        "Reports and KPIs",
    ]),
    ("Office Salespeople", "Desktop / Laptop", AMBER, [
        "Structured daily follow-up list",
        "Fast lead capture (<60 seconds)",
        "Qualification workflow",
        "Activity logging with reminders",
        "Prospect detail with full timeline",
    ]),
    ("Field Salespeople", "Mobile / On the Go", GREEN, [
        "15-second Quick Capture",
        "Pre-meeting client review",
        "Pipeline check between appointments",
        "Tap-to-call from follow-up cards",
        "Push notifications with quick actions",
    ]),
]

for i, (name, role, color, features) in enumerate(personas):
    x = Inches(1) + Inches(i * 3.9)
    card = add_shape(slide, x, Inches(1.8), Inches(3.5), Inches(4.8), WHITE, 0.03)
    card.line.color.rgb = color
    card.line.width = Pt(2)
    # Avatar circle
    avatar = add_shape(slide, x + Inches(1.25), Inches(2.1), Inches(1), Inches(1), color)
    # Make it a circle
    add_text_box(slide, x + Inches(1.25), Inches(2.3), Inches(1), Inches(0.6),
                 name[0], font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), Inches(3.2), Inches(2.9), Inches(0.4),
                 name, font_size=18, color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), Inches(3.6), Inches(2.9), Inches(0.3),
                 role, font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_bullet_list(slide, x + Inches(0.4), Inches(4.1), Inches(2.7), Inches(2.3),
                    [f"  {f}" for f in features], font_size=12, color=GRAY)


# ============================================================
# SLIDE 4: Pipeline Stages
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.7),
             "Pipeline Stages", font_size=36, color=TEAL_DARK, bold=True)
add_shape(slide, Inches(1), Inches(1.3), Inches(2), Inches(0.04), AMBER)

add_text_box(slide, Inches(1), Inches(1.6), Inches(11), Inches(0.5),
             "UK mortgage nurture process, designed around Mallard's specific pain points",
             font_size=16, color=GRAY)

stages = [
    ("New\nEnquiry", RGBColor(0xE0, 0xE7, 0xFF), RGBColor(0x43, 0x38, 0xCA), "5"),
    ("Initial\nContact", RGBColor(0xDB, 0xEA, 0xFE), RGBColor(0x1D, 0x4E, 0xD8), "3"),
    ("Not\nproceeded.", AMBER_LIGHT, RGBColor(0x92, 0x40, 0x0E), "12"),
    ("Nurturing", GREEN_LIGHT, RGBColor(0x16, 0x65, 0x34), "8"),
    ("DIP\ndone", RGBColor(0xCC, 0xFB, 0xF1), RGBColor(0x0F, 0x76, 0x6E), "2"),
    ("Ready to\nproceed", BLUE_LIGHT, RGBColor(0x1D, 0x4E, 0xD8), "4"),
    ("Referred\nto MAB", PURPLE_LIGHT, PURPLE, "6"),
    ("Won /\nCompleted", RGBColor(0xDC, 0xFC, 0xE7), GREEN, "3"),
    ("Lost /\nGone Cold", RGBColor(0xFE, 0xE2, 0xE2), RED, "9"),
]

for i, (name, bg_color, text_color, count) in enumerate(stages):
    x = Inches(0.7) + Inches(i * 1.55)
    y = Inches(2.5)
    card = add_shape(slide, x, y, Inches(1.35), Inches(1.8), bg_color, 0.05)
    if i == 2:  # Not proceeded. - highlighted
        card.line.color.rgb = AMBER
        card.line.width = Pt(3)
    else:
        card.line.color.rgb = text_color
        card.line.width = Pt(1)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.2), Inches(1.15), Inches(0.7),
                 name, font_size=12, color=text_color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(1.0), Inches(1.15), Inches(0.5),
                 count, font_size=28, color=text_color, bold=True, alignment=PP_ALIGN.CENTER)
    # Arrow between stages
    if i < len(stages) - 1:
        add_text_box(slide, x + Inches(1.35), y + Inches(0.65), Inches(0.2), Inches(0.4),
                     "→", font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

# Callout for "Not proceeded."
callout = add_shape(slide, Inches(2.8), Inches(4.7), Inches(7.5), Inches(1.2), AMBER_LIGHT, 0.03)
callout.line.color.rgb = AMBER
callout.line.width = Pt(2)
add_text_box(slide, Inches(3.1), Inches(4.85), Inches(7), Inches(0.35),
             '"Not proceeded." is a first-class stage — not an afterthought.',
             font_size=16, color=RGBColor(0x92, 0x40, 0x0E), bold=True)
add_text_box(slide, Inches(3.1), Inches(5.25), Inches(7), Inches(0.5),
             "This is where Mallard's leads currently fall through cracks. The system highlights it with amber borders, "
             "prominent follow-up scheduling, and supports 6-12 month reminder cycles.",
             font_size=13, color=RGBColor(0x92, 0x40, 0x0E))


# ============================================================
# Helper: Screen detail slide
# ============================================================
def make_screen_slide(title, subtitle, persona, device, node_id, features, color=TEAL_DARK):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)

    # Header bar
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), color)

    # Title + persona badge
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.6),
                 title, font_size=30, color=TEAL_DARK, bold=True)

    # Persona pill
    pill = add_shape(slide, Inches(0.8), Inches(1.05), Inches(1.6), Inches(0.35), color, 0.3)
    add_text_box(slide, Inches(0.8), Inches(1.05), Inches(1.6), Inches(0.35),
                 persona, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Device pill
    dev_color = BLUE if device == "Desktop" else GREEN
    pill2 = add_shape(slide, Inches(2.5), Inches(1.05), Inches(1.2), Inches(0.35), dev_color, 0.3)
    add_text_box(slide, Inches(2.5), Inches(1.05), Inches(1.2), Inches(0.35),
                 device, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.8), Inches(1.6), Inches(6), Inches(0.5),
                 subtitle, font_size=15, color=GRAY)

    # Features list
    add_text_box(slide, Inches(0.8), Inches(2.3), Inches(5), Inches(0.4),
                 "Key Features", font_size=18, color=BLACK, bold=True)
    add_bullet_list(slide, Inches(0.8), Inches(2.8), Inches(5.5), Inches(4),
                    [f"  {f}" for f in features], font_size=14, color=GRAY)

    # Screen placeholder (right side)
    if device == "Mobile":
        pw, ph = Inches(2.5), Inches(5.4)
        px = Inches(9.5)
        py = Inches(1.2)
        add_screen_placeholder(slide, px, py, pw, ph, f"Screen: {node_id}", "mobile")
    else:
        pw, ph = Inches(5.8), Inches(3.6)
        px = Inches(7)
        py = Inches(1.5)
        add_screen_placeholder(slide, px, py, pw, ph, f"Screen: {node_id}", "desktop")

    # Design file reference
    add_text_box(slide, Inches(7), Inches(6.5), Inches(5.5), Inches(0.5),
                 f"Design: pencil-welcome-desktop.pen  |  Node: {node_id}",
                 font_size=10, color=GRAY, alignment=PP_ALIGN.RIGHT)

    return slide


# ============================================================
# SLIDE 5: Manager Dashboard (Desktop)
# ============================================================
make_screen_slide(
    "Manager Dashboard", "Della sees everything at a glance — pipeline health, team performance, and which leads need attention.",
    "Della", "Desktop", "00PCT",
    [
        "4 KPI cards: New Leads, Follow-ups Due, Overdue, Deals Closed",
        "Pipeline Health horizontal bar chart across all active stages",
        "Team Activity live feed — calls logged, stages changed, leads created",
        "Real-time updates via Firestore onSnapshot",
        "Full sidebar navigation with manager-level sections",
        "In control within 30 seconds of opening the app",
    ]
)

# ============================================================
# SLIDE 6: Manager Dashboard (Mobile)
# ============================================================
make_screen_slide(
    "Manager Dashboard — Mobile", "Same data optimized for Della's phone. Check team status on the go.",
    "Della", "Mobile", "bO6Pj",
    [
        "Compact KPI pills: New leads, Overdue, Due today, Won",
        "Team Scorecard — avatar + stats per salesperson",
        "Condensed activity feed with color-coded dots",
        "Bottom tab nav with Dashboard active",
        "FAB button for quick actions",
    ]
)

# ============================================================
# SLIDE 7: Pipeline Board (Desktop Kanban)
# ============================================================
make_screen_slide(
    "Pipeline Board — Kanban", "Drag-and-drop Kanban board with 6 visible stage columns. Bird's-eye view of all prospects.",
    "Della", "Desktop", "yuS5D",
    [
        "6 color-coded stage columns: New Enquiry → Deal Done",
        '"Not proceeded." column highlighted with amber border — the key stage',
        "Prospect cards show: name, type, assigned avatar, time context",
        "Overdue badges in red on cards that need attention",
        "Filter and search in the toolbar",
        "47 total leads with per-stage count badges",
    ],
    color=AMBER
)

# ============================================================
# SLIDE 8: My Day Dashboard (Desktop)
# ============================================================
make_screen_slide(
    "\"My Day\" Dashboard", "The salesperson's home screen. Clear daily priorities with no ambiguity.",
    "Salesperson", "Desktop", "xLorR",
    [
        "Personalised greeting with motivational nudge",
        "Today's Follow-ups — cards with name, phone, context blurb, status badges",
        "Status badges: Overdue (red), Due today (amber), Planned (blue)",
        "My Pipeline summary — stage counts at a glance",
        "New Assignments — leads Della has assigned with context",
        "Focused on what to do right now",
    ]
)

# ============================================================
# SLIDE 9: My Day Dashboard (Mobile)
# ============================================================
make_screen_slide(
    "\"My Day\" — Mobile", "Single-column, follow-up focused. Everything works with one hand between appointments.",
    "Salesperson", "Mobile", "S14Wd",
    [
        "3 follow-up cards with avatar, name, phone number",
        "Status badges: Overdue, Due today, Planned",
        "Context blurb for each prospect",
        "Call and Snooze action buttons on overdue cards",
        "Bottom tab nav: My Day, Pipeline, +, Prospects, More",
        "FAB button for instant new lead capture",
    ]
)

# ============================================================
# SLIDE 10: New Lead Form (Desktop)
# ============================================================
make_screen_slide(
    "New Lead Form", "Two-column layout for capturing leads in under 60 seconds. The core data entry screen.",
    "Salesperson", "Desktop", "oluoT",
    [
        "Left: Contact details — name, phone, email, source, mortgage type, readiness, notes",
        "Right: Follow-up Reminder config — date, reason, up to 3 email recipients",
        "Up to 3 email recipients — ensures the right people get the reminder",
        "Reminder note — context to include in the email",
        "Smart defaults for follow-up dates based on readiness",
        "Save Lead button creates the lead + schedules the reminder in one action",
    ]
)

# ============================================================
# SLIDE 11: Quick Capture (Mobile)
# ============================================================
make_screen_slide(
    "Quick Capture", "15-second lead capture for networking events and in-person referrals. Minimal friction.",
    "Field Sales", "Mobile", "gJmkX",
    [
        "Only 2 required fields: Name and Phone",
        "Quick Note textarea for context",
        "Type tags — tap to select: First-time buyer, Remortgage, Self-employed, etc.",
        "Save button in header for one-tap completion",
        "Designed for capture while the conversation is still fresh",
        "Full details can be added later on desktop",
    ],
    color=GREEN
)

# ============================================================
# SLIDE 12: Prospect Detail (Desktop)
# ============================================================
make_screen_slide(
    "Prospect Detail", "Full prospect view with contact info, follow-up schedule, and activity timeline.",
    "Salesperson", "Desktop", "CHnFT",
    [
        "Header: avatar, name, tappable stage badge, meta info",
        "Action buttons: Call, Email, Log Activity",
        "Tabs: Overview, Notes & Activity, Qualification, Follow-ups",
        "Contact Information card with all key fields",
        "Next Follow-up card highlighted in amber — never missed",
        "Activity Timeline with color-coded dots and full context notes",
    ]
)

# ============================================================
# SLIDE 13: Pipeline List (Mobile)
# ============================================================
make_screen_slide(
    "Pipeline List — Mobile", "Mobile-optimized pipeline view. Horizontal stage pills for quick filtering.",
    "Field Sales", "Mobile", "xNuto",
    [
        "Horizontal stage pills — tap to filter by stage",
        "\"Not proceeded.\" selected, showing 12 leads",
        "Prospect cards with avatar, name, type summary",
        "Color-coded badges: Overdue, Due today, scheduled date, On track",
        "Bottom nav with Pipeline tab active",
        "Quick access to prospect detail via card tap",
    ],
    color=GREEN
)

# ============================================================
# SLIDE 14: MAB Import (Desktop)
# ============================================================
make_screen_slide(
    "MAB Platform Import", "Bridge to the existing system. CSV/XLS upload with automatic column mapping and deduplication.",
    "Della", "Desktop", "oaR0H",
    [
        "Drag-and-drop upload area with file confirmation",
        "Auto column mapping: MAB columns → Mallard Connect fields",
        "Deduplication preview with 3 color-coded buckets:",
        "    Green (42) — New leads, will be created",
        "    Amber (8) — Duplicates, will be skipped",
        "    Blue (3) — Duplicates with newer data, merge toggles",
        "Column mapper remembers last mapping",
        "Phone normalization handles +44 prefix and spacing",
    ],
    color=PURPLE
)

# ============================================================
# SLIDE 15: Core Feature - Follow-up Reminders
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), AMBER)
add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.7),
             "Core Feature: Follow-up Reminders", font_size=36, color=TEAL_DARK, bold=True)
add_shape(slide, Inches(1), Inches(1.2), Inches(2), Inches(0.04), AMBER)
add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.5),
             "The single most important feature. This is what Mallard Connect exists to do.",
             font_size=18, color=GRAY)

# Flow steps
steps = [
    ("1", "Set Date", "Smart defaults based\non readiness level", TEAL_DARK),
    ("2", "Add Recipients", "Up to 3 email addresses\nwho should be notified", AMBER),
    ("3", "Write Context", "Client situation blurb\nincluded in the email", BLUE),
    ("4", "7am Daily Cron", "Vercel Cron + Resend\nsends reminders", GREEN),
    ("5", "Act on Reminder", "Email with client details\n+ deep link to app", PURPLE),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = Inches(0.8) + Inches(i * 2.5)
    y = Inches(2.5)
    circle = add_shape(slide, x + Inches(0.65), y, Inches(0.8), Inches(0.8), color)
    add_text_box(slide, x + Inches(0.65), y + Inches(0.15), Inches(0.8), Inches(0.5),
                 num, font_size=24, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(1.0), Inches(2.1), Inches(0.4),
                 title, font_size=16, color=BLACK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(1.4), Inches(2.1), Inches(0.6),
                 desc, font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        add_text_box(slide, x + Inches(2.1), y + Inches(0.2), Inches(0.4), Inches(0.4),
                     "→", font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)

# Key insight box
insight = add_shape(slide, Inches(1), Inches(4.8), Inches(11.3), Inches(1.5), AMBER_LIGHT, 0.03)
insight.line.color.rgb = AMBER
insight.line.width = Pt(2)
add_text_box(slide, Inches(1.4), Inches(5.0), Inches(10.5), Inches(0.4),
             "Why this matters:", font_size=16, color=RGBColor(0x92, 0x40, 0x0E), bold=True)
add_text_box(slide, Inches(1.4), Inches(5.4), Inches(10.5), Inches(0.7),
             'The reminder email includes enough context that the salesperson can act without opening the app — '
             'client name, phone, situation summary, and a deep link to the prospect detail. '
             'Long nurture cycles (6-12 months) are normal in UK mortgage — the system handles them gracefully.',
             font_size=14, color=RGBColor(0x92, 0x40, 0x0E))


# ============================================================
# SLIDE 16: Tech Stack
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.7),
             "Tech Stack & Architecture", font_size=36, color=TEAL_DARK, bold=True)
add_shape(slide, Inches(1), Inches(1.3), Inches(2), Inches(0.04), AMBER)

stack_items = [
    ("Framework", "Next.js 14+ (App Router, TypeScript)", TEAL_DARK),
    ("Database", "Firebase Firestore (NoSQL, real-time sync)", AMBER),
    ("Auth", "Firebase Auth (email/password, magic link)", BLUE),
    ("Email", "Resend (transactional reminders, invites)", GREEN),
    ("Cron", "Vercel Cron (daily 7am UK reminder job)", PURPLE),
    ("UI", "Tailwind CSS + shadcn/ui + KokonutUI", TEAL_MID),
    ("Hosting", "Vercel (Next.js optimized, preview deploys)", BLUE),
    ("Validation", "Zod (shared client/server schemas)", GREEN),
    ("Real-time", "Firestore onSnapshot (no extra WebSocket infra)", AMBER),
    ("Import", "SheetJS/xlsx (CSV/XLS parsing)", PURPLE),
]

for i, (label, value, color) in enumerate(stack_items):
    col = i // 5
    row = i % 5
    x = Inches(1) + Inches(col * 6)
    y = Inches(1.8) + Inches(row * 0.85)
    dot = add_shape(slide, x, y + Inches(0.08), Inches(0.15), Inches(0.15), color)
    add_text_box(slide, x + Inches(0.3), y, Inches(1.5), Inches(0.3), label, font_size=14, color=BLACK, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(4.5), Inches(0.4), value, font_size=13, color=GRAY)

# Migration path
add_shape(slide, Inches(1), Inches(6.0), Inches(11.3), Inches(0.04), LIGHT_GRAY)
add_text_box(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
             "Migration path: Firebase → Supabase/Postgres when Mallard is ready for a full CRM",
             font_size=14, color=GRAY)


# ============================================================
# SLIDE 17: Phased Delivery
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.7),
             "Phased Delivery Plan", font_size=36, color=TEAL_DARK, bold=True)
add_shape(slide, Inches(1), Inches(1.3), Inches(2), Inches(0.04), AMBER)

phases = [
    ("Phase 1", "Foundation + Core", TEAL_DARK,
     "Lead intake, follow-up scheduler (3 email recipients), daily 7am cron job, basic dashboard, MAB CSV import, PWA"),
    ("Phase 2", "Pipeline + Visibility", AMBER,
     "Kanban board, stage change flow, manager dashboard, real-time activity feed, notification system"),
    ("Phase 3", "Team + Qualification", BLUE,
     "Team management, lead assignment, qualification fields, quick capture, meeting prep"),
    ("Phase 4", "Reports + Polish", GREEN,
     "KPI reports, leaderboard, export, onboarding wizard, edge cases, audit log"),
]

for i, (phase, title, color, desc) in enumerate(phases):
    x = Inches(0.8) + Inches(i * 3.15)
    card = add_shape(slide, x, Inches(1.8), Inches(2.9), Inches(4.5), WHITE, 0.03)
    card.line.color.rgb = color
    card.line.width = Pt(2)
    # Phase header
    hdr = add_shape(slide, x, Inches(1.8), Inches(2.9), Inches(0.8), color, 0.03)
    add_text_box(slide, x + Inches(0.2), Inches(1.9), Inches(2.5), Inches(0.3),
                 phase, font_size=12, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.2), Inches(2.2), Inches(2.5), Inches(0.3),
                 title, font_size=16, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.2), Inches(2.8), Inches(2.5), Inches(3.3),
                 desc, font_size=13, color=GRAY)

# Future
add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
             "Future: CRM Graduation → Supabase/Postgres, email integration, calendar sync, automated workflows",
             font_size=13, color=GRAY)


# ============================================================
# SLIDE 18: Thank You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, TEAL_DARK)

add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1),
             "MALLARD CONNECT", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(5.5), Inches(3.2), Inches(2.3), Inches(0.04), AMBER)
add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.6),
             "No prospect falls through the cracks.", font_size=22, color=AMBER, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
             "Designed for Mallard Mortgages  |  Built by Storyboard Digital  |  Sheffield, UK",
             font_size=14, color=RGBColor(0x7D, 0xB3, 0xB0), alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
             "UI Mockups: pencil-welcome-desktop.pen  |  10 screens across desktop & mobile",
             font_size=12, color=RGBColor(0x5A, 0x8A, 0x88), alignment=PP_ALIGN.CENTER)


# ============================================================
# Save
# ============================================================
output_dir = "/Users/vinn/conductor/workspaces/mallard_connect/sao-paulo/docs"
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, "Mallard-Connect-UX-Presentation.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
